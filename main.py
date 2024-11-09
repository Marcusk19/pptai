import os
from pptx import Presentation
from pptx.util import Inches
from client import setup_client
from images import generate_image, download_image

import yaml

# Set up your OpenAI client
client = setup_client("config.yml")


def generate_content(prompt):
    # Generate content using OpenAI API (GPT model)
    try:
        response = client.chat.completions.create(
            model="gpt-3.5-turbo",  # Using the newer GPT-3.5-turbo model
            messages=[{"role": "user", "content": prompt}],
            max_tokens=500,  # You may adjust tokens based on the content length
        )
        return response.choices[0].message.content.strip()
    except Exception as e:
        print(f"Error generating content: {str(e)}")
        return None


def create_ppt(yaml_content):
    # Parse YAML content
    try:
        slides_data = yaml.safe_load(yaml_content)
    except yaml.YAMLError as e:
        print(f"Error parsing YAML content: {str(e)}")
        return

    # Initialize a PowerPoint presentation
    prs = Presentation()

    for slide in slides_data:
        # Each slide is a dictionary with one key (slide1, slide2, etc.)
        slide_content = list(slide.values())[0]

        # Add a new slide with title and content layout
        ppt_slide = prs.slides.add_slide(prs.slide_layouts[1])

        # Set the title
        title = slide_content.get("title", "")
        ppt_slide.shapes.title.text = title

        # Get the content placeholder
        content_placeholder = ppt_slide.shapes.placeholders[1]

        # Add a image
        try:
            image_prompt = slide_content.get("image_prompt", "a missing icon")
            image_url = generate_image(image_prompt)
            image = download_image(image_url)
            image_path = f"{title}.png"
            image.save(image_path)

            # Calculate dimensions
            text_width = Inches(5)  # Width for text
            image_width = Inches(4)  # Width for image

            # Left side for text content (keep original placeholder)
            content_placeholder.width = text_width
            content_placeholder.height = Inches(
                4.5
            )  # Slightly shorter to avoid overlap
            content_placeholder.top = Inches(2)  # Move content down below title

            image_left = (
                content_placeholder.left + text_width + Inches(0.5)
            )  # Add 0.5 inch padding
            image_top = Inches(2)  # Align with content top

            ppt_slide.shapes.add_picture(
                image_path, image_left, image_top, width=image_width
            )
            os.remove(image_path)
        except Exception as e:
            print(f"Error adding image: {str(e)}")

        text_frame = content_placeholder.text_frame
        text_frame.clear()  # Clear existing text

        # Add each bullet point
        points = slide_content.get("points", [])
        for point in points:
            p = text_frame.add_paragraph()
            p.text = point
            p.level = 0  # 0 for main points, increase for sub-points

        print("Slide added:", title)

        # Add notes
        # notes_slide = ppt_slide.notes_slide
        # text_frame = notes_slide.notes_text_frame
        # text_frame.text = slide_content.get('script', '')

    # Save the presentation
    prs.save("generated_presentation.pptx")
    print("Presentation saved as 'generated_presentation.pptx'")


# Main script
if __name__ == "__main__":
    prompt = input("Enter your prompt for the presentation content: ")
    rich_prompt = f"""
        {prompt}
        Create a presentation outline for this.
        Format the presentation and script in yaml so that each
        slide is it's own block with points underneath it that will be part of the slide.
        Number the slides in the yaml. Also add title fields for each slide.
        Finally, add an image_prompt field that contains a detailed prompt that can be fed to
        an AI model to generate a helpful image for the slide.

        The format should be formatted like this:

        - slide1:
            title: "title"
            points:
              - "point1"
              - "point2"
              - "point3"
            image_prompt: "image prompt to be fed to AI model"

        All fields should be properly sanitized and formatted. There should be no special characters
        in the fields like colons.
        """
    generated_text = generate_content(rich_prompt)
    print(generated_text)

    if generated_text:
        yaml_content = generated_text
        # parse yaml delimiters if present
        if "```yaml" in generated_text:
            yaml_content = generated_text.split("```yaml")[-1]
            yaml_content = yaml_content.split("```")[0]
            yaml_content = yaml_content.strip()
        create_ppt(generated_text)
        pass
    else:
        print("Failed to generate presentation content")
